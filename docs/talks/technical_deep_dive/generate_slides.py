"""
Generate technical_deep_dive PPTX slides for DASH_GUI talk.
Run from repo root:  uv run python docs/talks/technical_deep_dive/generate_slides.py
Output: docs/talks/technical_deep_dive/technical_deep_dive_slides.pptx
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# ── Colour palette ────────────────────────────────────────────────────────────
C_BG      = RGBColor(0x1E, 0x29, 0x3B)
C_ACCENT  = RGBColor(0x38, 0xBD, 0xF8)
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT   = RGBColor(0xC8, 0xD6, 0xE8)
C_GREEN   = RGBColor(0x4A, 0xDE, 0x80)
C_YELLOW  = RGBColor(0xFB, 0xBF, 0x24)
C_RED     = RGBColor(0xF8, 0x71, 0x71)
C_CODE_BG = RGBColor(0x0F, 0x17, 0x2A)
C_ORANGE  = RGBColor(0xFB, 0x92, 0x3C)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def fill_bg(slide, color=C_BG):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def add_textbox(slide, text, left, top, width, height,
                font_size=20, bold=False, color=C_WHITE,
                align=PP_ALIGN.LEFT, word_wrap=True):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(font_size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tb


def add_para(tf, text, font_size=16, bold=False, color=C_WHITE, align=PP_ALIGN.LEFT):
    p = tf.add_paragraph()
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(font_size)
    r.font.bold = bold
    r.font.color.rgb = color
    return p


def section_pill(slide, text, left=Inches(0.4), top=Inches(0.25)):
    pill = slide.shapes.add_shape(9, left, top, Inches(2.6), Inches(0.38))
    pill.fill.solid()
    pill.fill.fore_color.rgb = C_ACCENT
    pill.line.fill.background()
    tf = pill.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(13)
    r.font.bold = True
    r.font.color.rgb = C_BG


def heading(slide, text, font_size=34):
    add_textbox(slide, text, Inches(0.4), Inches(0.75), Inches(12.5), Inches(0.9),
                font_size=font_size, bold=True, color=C_WHITE)


def add_code_box(slide, code, left, top, width, height, font_size=13, border_color=C_ACCENT):
    box = slide.shapes.add_shape(1, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = C_CODE_BG
    box.line.color.rgb = border_color
    box.line.width = Pt(1)
    tf = box.text_frame
    tf.word_wrap = False
    first = True
    for line in code.split("\n"):
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = line
        r.font.size = Pt(font_size)
        r.font.color.rgb = C_ACCENT
        r.font.name = "Courier New"


def accent_bar(slide, color=C_ACCENT):
    bar = slide.shapes.add_shape(1, 0, 0, Inches(0.18), SLIDE_H)
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()


# ── Slides ────────────────────────────────────────────────────────────────────

def slide_title(prs):
    sl = blank_slide(prs)
    fill_bg(sl)
    accent_bar(sl)
    add_textbox(sl, "DASH_GUI: A Technical Deep Dive",
                Inches(0.5), Inches(1.4), Inches(12.5), Inches(1.6),
                font_size=48, bold=True, color=C_WHITE)
    add_textbox(sl, "Architecture · Enrichment Engine · SHACL Generation",
                Inches(0.5), Inches(3.1), Inches(12), Inches(0.8),
                font_size=24, color=C_ACCENT)
    add_textbox(sl, "Audience: Developers, architects, semantic web practitioners",
                Inches(0.5), Inches(4.1), Inches(10), Inches(0.6),
                font_size=18, color=C_LIGHT)
    add_textbox(sl, "45–60 minutes",
                Inches(0.5), Inches(5.5), Inches(4), Inches(0.5),
                font_size=16, color=C_YELLOW)


def slide_agenda(prs):
    sl = blank_slide(prs)
    fill_bg(sl)
    section_pill(sl, "Agenda")
    heading(sl, "What We'll Cover")
    items = [
        ("1", "The Brick Approach",      "Semantic component reuse"),
        ("2", "Architecture Overview",   "Multi-tenant, shared core"),
        ("3", "The Enrichment Engine",   "4-layer resolution — generic, not QUDT-only"),
        ("4", "SHACL Export Pipeline",   "From bricks to validation rules"),
        ("5", "Web vs Qt Frontends",     "Same core, different UI technology"),
        ("6", "Extension Points",        "Add ontologies, rules, templates"),
    ]
    top = Inches(1.9)
    for num, title, desc in items:
        circle = sl.shapes.add_shape(9, Inches(0.4), top, Inches(0.5), Inches(0.5))
        circle.fill.solid(); circle.fill.fore_color.rgb = C_ACCENT
        circle.line.fill.background()
        ctf = circle.text_frame
        cp = ctf.paragraphs[0]; cp.alignment = PP_ALIGN.CENTER
        cr = cp.add_run(); cr.text = num; cr.font.size = Pt(15)
        cr.font.bold = True; cr.font.color.rgb = C_BG
        add_textbox(sl, title, Inches(1.1), top + Inches(0.03), Inches(3.5), Inches(0.48),
                    font_size=18, bold=True, color=C_WHITE)
        add_textbox(sl, desc,  Inches(4.8), top + Inches(0.06), Inches(8.2), Inches(0.42),
                    font_size=16, color=C_LIGHT)
        top += Inches(0.82)


def slide_why_shacl(prs):
    sl = blank_slide(prs)
    fill_bg(sl)
    section_pill(sl, "Motivation")
    heading(sl, "Why Traditional SHACL Authoring Fails")

    left_box = sl.shapes.add_shape(1, Inches(0.4), Inches(1.9), Inches(5.8), Inches(4.5))
    left_box.fill.solid(); left_box.fill.fore_color.rgb = RGBColor(0x3B,0x08,0x08)
    left_box.line.color.rgb = C_RED; left_box.line.width = Pt(1.5)
    tf = left_box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "Traditional Approach"; r.font.size = Pt(17)
    r.font.bold = True; r.font.color.rgb = C_RED
    for line in ["• Write raw Turtle by hand",
                 "• Manual unit lists (copy-paste)",
                 "• No reuse across schemas",
                 "• Validation errors found late",
                 "• Ontology experts only"]:
        add_para(tf, line, font_size=16, color=C_LIGHT)

    right_box = sl.shapes.add_shape(1, Inches(6.8), Inches(1.9), Inches(5.8), Inches(4.5))
    right_box.fill.solid(); right_box.fill.fore_color.rgb = RGBColor(0x05,0x2E,0x16)
    right_box.line.color.rgb = C_GREEN; right_box.line.width = Pt(1.5)
    tf2 = right_box.text_frame; tf2.word_wrap = True
    p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = "Brick Approach"; r2.font.size = Pt(17)
    r2.font.bold = True; r2.font.color.rgb = C_GREEN
    for line in ["• Compose visual bricks",
                 "• Ontology-driven widgets (auto)",
                 "• Reuse across all schemas",
                 "• Enrichment hints while authoring",
                 "• Domain experts can author"]:
        add_para(tf2, line, font_size=16, color=C_LIGHT)


def slide_brick_anatomy(prs):
    sl = blank_slide(prs)
    fill_bg(sl)
    section_pill(sl, "The Brick Approach")
    heading(sl, "What Is a Brick?")

    add_textbox(sl, "A Brick = NodeShape + PropertyShapes, stored as JSON, exported as SHACL.",
                Inches(0.4), Inches(1.85), Inches(12.5), Inches(0.55),
                font_size=18, color=C_LIGHT)

    code = (
        '{\n'
        '  "name": "Mass",\n'
        '  "object_type": "NodeShape",\n'
        '  "target_class": "ex:Mass",\n'
        '  "leaf_properties": [{\n'
        '    "path": "ex:value",\n'
        '    "datatype": "xsd:decimal",\n'
        '    "sh_class": "qudt:Mass",   ← triggers enrichment\n'
        '    "min_count": 1,\n'
        '    "max_count": 1\n'
        '  }]\n'
        '}'
    )
    add_code_box(sl, code, Inches(0.4), Inches(2.55), Inches(6.0), Inches(4.2), font_size=15)

    notes = [
        ("sh:class",    "key field — drives widget selection"),
        ("min/maxCount","cardinality constraints"),
        ("JSON storage","fast CRUD, git-friendly"),
        ("Export",      "→ SHACL Turtle + dash:editor annotations"),
    ]
    top = Inches(2.55)
    for field, desc in notes:
        add_textbox(sl, field, Inches(6.7), top, Inches(2.4), Inches(0.48),
                    font_size=16, bold=True, color=C_ACCENT, align=PP_ALIGN.RIGHT)
        add_textbox(sl, desc,  Inches(9.3), top, Inches(3.8), Inches(0.48),
                    font_size=15, color=C_LIGHT)
        top += Inches(0.9)


def slide_architecture(prs):
    sl = blank_slide(prs)
    fill_bg(sl)
    section_pill(sl, "Architecture")
    heading(sl, "System Architecture")

    layers = [
        ("Frontend Layer",   "Web (React/Babel) :5001  ·  Qt (PyQt6) desktop", C_ACCENT,  Inches(1.85)),
        ("REST API",         "Flask — multi-tenant, per-session BrickCore",     C_LIGHT,   Inches(2.75)),
        ("Core Services",    "OntologyManager · EnrichmentEngine · SHACLExporter", C_YELLOW, Inches(3.65)),
        ("Shared Library",   "JSON bricks/schemas  ·  TTL ontology cache",     C_GREEN,   Inches(4.55)),
        ("Standards Out",    "SHACL Turtle  ·  DASH editors  ·  HTML forms",   C_ORANGE,  Inches(5.45)),
    ]
    for label, desc, col, top in layers:
        box = sl.shapes.add_shape(1, Inches(1.2), top, Inches(10.8), Inches(0.72))
        box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0x0F,0x22,0x36)
        box.line.color.rgb = col; box.line.width = Pt(2)
        tf = box.text_frame; tf.word_wrap = False
        p = tf.paragraphs[0]
        r1 = p.add_run(); r1.text = f"{label:<22}"; r1.font.size = Pt(17)
        r1.font.bold = True; r1.font.color.rgb = col; r1.font.name = "Courier New"
        r2 = p.add_run(); r2.text = desc; r2.font.size = Pt(15)
        r2.font.color.rgb = C_LIGHT; r2.font.name = "Courier New"

    add_textbox(sl, "Both Web and Qt share the same BrickCore — no logic duplication.",
                Inches(0.4), Inches(6.5), Inches(12.5), Inches(0.5),
                font_size=16, color=C_YELLOW)


def slide_shared_library(prs):
    sl = blank_slide(prs)
    fill_bg(sl)
    section_pill(sl, "Shared Library")
    heading(sl, "The Shared Library System")

    tree = (
        "shared_libraries/\n"
        "├── bricks/default/bricks/\n"
        "│   ├── mass_nodeshape.json\n"
        "│   ├── temperature_nodeshape.json\n"
        "│   └── ...\n"
        "├── schemas/default/schemas/\n"
        "│   └── ev_battery_schema.json\n"
        "└── ontologies/cache/\n"
        "    ├── qudt-units.ttl       (2906 units)\n"
        "    ├── qudt_quantitykind.ttl (1217 kinds)\n"
        "    ├── skos.rdf\n"
        "    ├── foaf.ttl\n"
        "    └── schema.org.ttl       (+ 14 more)"
    )
    add_code_box(sl, tree, Inches(0.4), Inches(1.9), Inches(6.2), Inches(4.9), font_size=14)

    points = [
        "Filesystem-based — git-friendly, no database",
        "Single source of truth for both Web and Qt",
        "17 ontologies loaded at startup (~45 MB TTL)",
        "Add ontology → drop .ttl in cache → restart",
        "New prefixes & enrichment rules available instantly",
    ]
    top = Inches(2.1)
    for pt in points:
        add_textbox(sl, f"→  {pt}", Inches(6.8), top, Inches(6.1), Inches(0.52),
                    font_size=15, color=C_LIGHT)
        top += Inches(0.72)


def slide_enrichment_layer0(prs):
    sl = blank_slide(prs)
    fill_bg(sl)
    section_pill(sl, "Enrichment Engine")
    heading(sl, "Layer 0: Datatype → Widget  (fast path)")

    mapping = [
        ("xsd:boolean",    "boolean_toggle",  "checkbox"),
        ("xsd:date",       "date_picker",     "calendar widget"),
        ("xsd:decimal",    "decimal_input",   "numeric field"),
        ("xsd:anyURI",     "uri_input",       "URL-validated field"),
        ("rdf:langString", "language_text",   "multilingual input"),
        ("xsd:integer",    "decimal_input",   "integer-constrained"),
    ]
    headers = ("sh:datatype", "widget type", "rendered as")
    col_l = [Inches(0.5), Inches(4.8), Inches(8.5)]
    top = Inches(2.0)
    for i, h in enumerate(headers):
        add_textbox(sl, h, col_l[i], top, Inches(4.0), Inches(0.45),
                    font_size=16, bold=True, color=C_ACCENT)
    top += Inches(0.5)
    for dtype, widget, render in mapping:
        add_textbox(sl, dtype,  col_l[0], top, Inches(4.1), Inches(0.45),
                    font_size=15, color=C_WHITE)
        add_textbox(sl, widget, col_l[1], top, Inches(3.5), Inches(0.45),
                    font_size=15, color=C_GREEN)
        add_textbox(sl, render, col_l[2], top, Inches(4.5), Inches(0.45),
                    font_size=15, color=C_LIGHT)
        top += Inches(0.62)

    add_textbox(sl, "No ontology lookup — purely declarative rules in widget_rules.ttl. Add new datatypes by adding TTL triples.",
                Inches(0.4), Inches(6.5), Inches(12.5), Inches(0.55),
                font_size=15, color=C_YELLOW)


def slide_enrichment_layers12(prs):
    sl = blank_slide(prs)
    fill_bg(sl)
    section_pill(sl, "Enrichment Engine")
    heading(sl, "Layers 1-2: Semantic Resolution  (QUDT is one example)")

    # left panel — QUDT
    left = sl.shapes.add_shape(1, Inches(0.4), Inches(1.9), Inches(5.9), Inches(4.8))
    left.fill.solid(); left.fill.fore_color.rgb = RGBColor(0x0A,0x1E,0x30)
    left.line.color.rgb = C_ACCENT; left.line.width = Pt(1.5)
    tf = left.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "Physical quantities (QUDT)"; r.font.size = Pt(16)
    r.font.bold = True; r.font.color.rgb = C_ACCENT
    qudt_lines = [
        "sh:class = qudt:Mass",
        "  ↓ Layer 1: dimensional sig",
        "    (1, 0, 0, 0, 0, 0, 0)  [M¹]",
        "  ↓ Layer 2: applicableUnit",
        "    query QUDT graph",
        "",
        "→ unit_dropdown",
        "  [kg, g, lb, oz, mg, t]",
    ]
    for line in qudt_lines:
        add_para(tf, line, font_size=14, color=C_LIGHT)

    # right panel — SKOS
    right = sl.shapes.add_shape(1, Inches(6.9), Inches(1.9), Inches(5.9), Inches(4.8))
    right.fill.solid(); right.fill.fore_color.rgb = RGBColor(0x0A,0x1E,0x30)
    right.line.color.rgb = C_GREEN; right.line.width = Pt(1.5)
    tf2 = right.text_frame; tf2.word_wrap = True
    p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = "Controlled vocab (SKOS)"; r2.font.size = Pt(16)
    r2.font.bold = True; r2.font.color.rgb = C_GREEN
    skos_lines = [
        "sh:class = skos:Concept",
        "+ skos:inScheme <myVocab>",
        "  ↓ Layer 2: extract concepts",
        "    from scheme graph",
        "",
        "→ skos_selector dropdown",
        "  [Concept A, B, C ...]",
        "  (your vocabulary, not QUDT)",
    ]
    for line in skos_lines:
        add_para(tf2, line, font_size=14, color=C_LIGHT)

    add_textbox(sl, "Any ontology with the right predicate structure works — rules declared in widget_rules.ttl, zero Python changes.",
                Inches(0.4), Inches(6.85), Inches(12.5), Inches(0.45),
                font_size=14, color=C_YELLOW)


def slide_enrichment_layers34(prs):
    sl = blank_slide(prs)
    fill_bg(sl)
    section_pill(sl, "Enrichment Engine")
    heading(sl, "Layers 3-4: Namespace & Inheritance Fallbacks")

    table_data = [
        ("Layer", "Trigger",                   "Result",                        "Example",                    C_ACCENT),
        ("3",     "Namespace prefix match",     "Suggest known properties",      "foaf:Person → name, mbox…",  C_LIGHT),
        ("3",     "Predicate in loaded onto",   "Suggest known properties",      "schema:Thing → identifier…", C_LIGHT),
        ("4",     "rdfs:subClassOf chain",      "Entity search-as-you-type",     "foaf:Person ⊑ foaf:Agent",   C_LIGHT),
        ("—",     "No match at any layer",      "Plain text field (fallback)",   "unknown:Widget → text input",C_LIGHT),
    ]
    col_l   = [Inches(0.4), Inches(1.3), Inches(4.5), Inches(7.8)]
    col_w   = [Inches(0.7), Inches(3.0), Inches(3.1), Inches(5.0)]
    top = Inches(1.95)
    for row in table_data:
        *cells, row_color = row
        for i, cell in enumerate(cells):
            add_textbox(sl, cell, col_l[i], top, col_w[i], Inches(0.52),
                        font_size=15,
                        bold=(top == Inches(1.95)),
                        color=C_ACCENT if top == Inches(1.95) else row_color)
        top += Inches(0.72)

    add_textbox(sl, "Load any OWL/RDFS ontology → its classes and subclass chains become immediately available.",
                Inches(0.4), Inches(5.9), Inches(12.5), Inches(0.5),
                font_size=15, color=C_YELLOW)
    add_textbox(sl, "No vocabulary names are hardcoded in Python — purely data-driven.",
                Inches(0.4), Inches(6.45), Inches(12.5), Inches(0.5),
                font_size=15, color=C_GREEN)


def slide_widget_rules(prs):
    sl = blank_slide(prs)
    fill_bg(sl)
    section_pill(sl, "Declarative Config")
    heading(sl, "widget_rules.ttl — Rules Without Code Changes")

    code = (
        "# Physical quantity — dimensional signature\n"
        "dashgui:MassRule a dashgui:WidgetRule ;\n"
        "    dashgui:widget              \"unit_dropdown\" ;\n"
        "    dashgui:dimensionalSignature \"1 0 0 0 0 0 0\" ;\n"
        "    dashgui:siUnit <qudt-unit:KiloGM> .\n"
        "\n"
        "# Controlled vocabulary — predicate trigger\n"
        "dashgui:SkosRule a dashgui:WidgetRule ;\n"
        "    dashgui:widget          \"skos_selector\" ;\n"
        "    dashgui:triggerPredicate skos:inScheme .\n"
        "\n"
        "# Entity hierarchy — subClassOf trigger\n"
        "dashgui:AgentRule a dashgui:WidgetRule ;\n"
        "    dashgui:widget          \"entity_lookup\" ;\n"
        "    dashgui:triggerSubClassOf foaf:Agent ."
    )
    add_code_box(sl, code, Inches(0.4), Inches(1.85), Inches(8.5), Inches(4.9), font_size=14)

    points = [
        "No ontology names in Python",
        "12 dimensional signatures",
        "4 predicate-trigger rules",
        "3 subClassOf-trigger rules",
        "Add your rule → add TTL triples",
    ]
    top = Inches(2.2)
    for pt in points:
        add_textbox(sl, f"✓  {pt}", Inches(9.2), top, Inches(3.9), Inches(0.5),
                    font_size=15, color=C_GREEN)
        top += Inches(0.78)


def slide_shacl_export(prs):
    sl = blank_slide(prs)
    fill_bg(sl)
    section_pill(sl, "SHACL Export")
    heading(sl, "From Bricks to SHACL Turtle")

    left_code = (
        "# Brick (JSON)\n"
        "{\n"
        " \"name\": \"Mass\",\n"
        " \"object_type\": \"NodeShape\",\n"
        " \"leaf_properties\": [{\n"
        "   \"path\": \"ex:value\",\n"
        "   \"datatype\": \"xsd:decimal\",\n"
        "   \"sh_class\": \"qudt:Mass\"\n"
        " }]\n"
        "}"
    )
    right_code = (
        "# SHACL (Turtle)\n"
        "schema:Mass a sh:NodeShape ;\n"
        "  sh:targetClass ex:Mass ;\n"
        "  sh:property [\n"
        "    sh:path ex:value ;\n"
        "    sh:datatype xsd:decimal ;\n"
        "    sh:class qudt:Mass ;\n"
        "    dash:editor\n"
        "      dash:InstancesSelectEditor ;\n"
        "    sh:in (unit:KG unit:G\n"
        "           unit:LB unit:OZ) ;\n"
        "  ] ."
    )
    add_code_box(sl, left_code,  Inches(0.4), Inches(1.9), Inches(5.7), Inches(4.5), font_size=14)
    add_textbox(sl, "export_schema()\n       →", Inches(6.3), Inches(3.6), Inches(1.5), Inches(0.9),
                font_size=18, bold=True, color=C_YELLOW, align=PP_ALIGN.CENTER)
    add_code_box(sl, right_code, Inches(7.5), Inches(1.9), Inches(5.4), Inches(4.5), font_size=14)

    add_textbox(sl, "EnrichmentEngine runs at export time — dash:editor and sh:in resolved from sh:class via ontology lookup.",
                Inches(0.4), Inches(6.6), Inches(12.5), Inches(0.5),
                font_size=14, color=C_YELLOW)


def slide_dash_mappings(prs):
    sl = blank_slide(prs)
    fill_bg(sl)
    section_pill(sl, "DASH Mappings")
    heading(sl, "Widget Types → DASH Editors")

    rows = [
        ("text",           "dash:TextFieldEditor",        "dash:LabelViewer"),
        ("boolean_toggle", "dash:BooleanSelectEditor",    "dash:BooleanViewer"),
        ("date_picker",    "dash:DatePickerEditor",       "dash:LabelViewer"),
        ("decimal_input",  "dash:DecimalFieldEditor",     "dash:LabelViewer"),
        ("unit_dropdown",  "dash:InstancesSelectEditor",  "dash:LabelViewer"),
        ("skos_selector",  "dash:InstancesSelectEditor",  "dash:LabelViewer"),
        ("entity_lookup",  "dash:AutoCompleteEditor",     "dash:LabelViewer"),
    ]
    headers = ("widget type", "dash:editor", "dash:viewer")
    col_l = [Inches(0.4), Inches(4.3), Inches(9.0)]
    col_w = [Inches(3.7), Inches(4.5), Inches(4.0)]
    top = Inches(1.9)
    for i, h in enumerate(headers):
        add_textbox(sl, h, col_l[i], top, col_w[i], Inches(0.45),
                    font_size=16, bold=True, color=C_ACCENT)
    top += Inches(0.5)
    for widget, editor, viewer in rows:
        add_textbox(sl, widget, col_l[0], top, col_w[0], Inches(0.48),
                    font_size=14, color=C_WHITE)
        add_textbox(sl, editor, col_l[1], top, col_w[1], Inches(0.48),
                    font_size=14, color=C_GREEN)
        add_textbox(sl, viewer, col_l[2], top, col_w[2], Inches(0.48),
                    font_size=14, color=C_LIGHT)
        top += Inches(0.6)

    add_textbox(sl, "DASH (Data Shapes) provides the form widgets — SHACL provides the validation constraints.",
                Inches(0.4), Inches(6.55), Inches(12.5), Inches(0.5),
                font_size=14, color=C_YELLOW)


def slide_web_frontend(prs):
    sl = blank_slide(prs)
    fill_bg(sl)
    section_pill(sl, "Frontend: Web")
    heading(sl, "Web Frontend: React + Babel  (single file)")

    tree = (
        "index.html  (inline JSX, Babel standalone — no build step)\n"
        "└── App\n"
        "    ├── SessionManager\n"
        "    ├── BrickList  (sidebar)\n"
        "    └── BrickEditor\n"
        "        ├── BasicFields\n"
        "        ├── LeafPropertiesList\n"
        "        └── PropertyEditorModal\n"
        "            ├── DatatypeSelector\n"
        "            ├── OntologyBrowser  (all 17 ontologies)\n"
        "            └── EnrichmentDisplay  (live widget preview)"
    )
    add_code_box(sl, tree, Inches(0.4), Inches(1.9), Inches(7.8), Inches(4.5), font_size=14)

    points = [
        ("REST API",     "Flask backend, per-session BrickCore"),
        ("Enrichment",   "useEffect hooks call /api/enrichment"),
        ("Ontology",     "Browser opens on 'All Ontologies' merged view"),
        ("No build",     "Babel transpiles in-browser — rapid iteration"),
        ("Future",       "Migrate to Vite + React for production DX"),
    ]
    top = Inches(2.1)
    for label, desc in points:
        add_textbox(sl, label, Inches(8.5),  top, Inches(1.8), Inches(0.48),
                    font_size=15, bold=True, color=C_ACCENT)
        add_textbox(sl, desc,  Inches(10.45), top, Inches(2.7), Inches(0.48),
                    font_size=13, color=C_LIGHT)
        top += Inches(0.78)


def slide_qt_frontend(prs):
    sl = blank_slide(prs)
    fill_bg(sl)
    section_pill(sl, "Frontend: Qt")
    heading(sl, "Qt Frontend: PyQt6 + UI Files")

    tree = (
        "brick_editor.py  (main window)\n"
        "├── loadUi('property_editor.ui')   ← Qt Designer\n"
        "│   ├── Constraint group boxes (dynamic show/hide)\n"
        "│   ├── Datatype combo → _update_group_visibility()\n"
        "│   └── Class browser button → SimpleOntologyBrowser\n"
        "├── gui_components.py\n"
        "│   ├── PropertyEditorDialog\n"
        "│   ├── ConstraintEditorDialog\n"
        "│   └── SimpleOntologyBrowser\n"
        "└── brick_core_simple.py  ← same as web backend"
    )
    add_code_box(sl, tree, Inches(0.4), Inches(1.9), Inches(8.2), Inches(4.5), font_size=14)

    points = [
        ("UI files",   ".ui XML — layout changes without recompile"),
        ("loadUi()",   "Dynamic load at runtime — no pyuic step"),
        ("Same core",  "BrickCore shared with Flask web API"),
        ("Parity",     "Same constraint fields as web modal"),
    ]
    top = Inches(2.3)
    for label, desc in points:
        add_textbox(sl, label, Inches(8.9),  top, Inches(1.8), Inches(0.48),
                    font_size=15, bold=True, color=C_ACCENT)
        add_textbox(sl, desc,  Inches(10.85), top, Inches(2.3), Inches(0.48),
                    font_size=13, color=C_LIGHT)
        top += Inches(0.82)

    add_textbox(sl, "Key Point: Same BrickCore, different UI technology — no logic duplication.",
                Inches(0.4), Inches(6.55), Inches(12.5), Inches(0.5),
                font_size=15, color=C_YELLOW)


def slide_extension(prs):
    sl = blank_slide(prs)
    fill_bg(sl)
    section_pill(sl, "Extension Points")
    heading(sl, "How to Extend the System")

    rows = [
        ("Add ontology",      "Drop .ttl in ontologies/cache/",   "New classes & prefixes available"),
        ("Add widget rule",   "Add triples to widget_rules.ttl",  "New trigger → new widget, no Python"),
        ("Add datatype",      "Add rule to widget_rules.ttl",     "New datatype → new widget type"),
        ("Custom validator",  "Subclass BrickCore",               "Hook into save/load pipeline"),
        ("Custom exporter",   "Subclass SHACLExporter",           "Custom Turtle templates or formats"),
    ]
    headers = ("What",  "How",  "Result")
    col_l = [Inches(0.4), Inches(3.9), Inches(8.3)]
    col_w = [Inches(3.3), Inches(4.2), Inches(4.8)]
    top = Inches(1.9)
    for i, h in enumerate(headers):
        add_textbox(sl, h, col_l[i], top, col_w[i], Inches(0.45),
                    font_size=16, bold=True, color=C_ACCENT)
    top += Inches(0.5)
    for what, how, result in rows:
        add_textbox(sl, what,   col_l[0], top, col_w[0], Inches(0.52),
                    font_size=15, bold=True, color=C_WHITE)
        add_textbox(sl, how,    col_l[1], top, col_w[1], Inches(0.52),
                    font_size=14, color=C_LIGHT)
        add_textbox(sl, result, col_l[2], top, col_w[2], Inches(0.52),
                    font_size=14, color=C_GREEN)
        top += Inches(0.72)

    add_textbox(sl, "Everything is pluggable — no ontology names hardcoded in Python.",
                Inches(0.4), Inches(6.55), Inches(12.5), Inches(0.5),
                font_size=15, color=C_YELLOW)


def slide_numbers(prs):
    sl = blank_slide(prs)
    fill_bg(sl)
    section_pill(sl, "Performance")
    heading(sl, "Numbers")

    metrics = [
        ("Ontologies loaded at startup",  "17  (≈ 45 MB TTL)"),
        ("QUDT units indexed",            "2,906"),
        ("Quantity kinds indexed",        "1,217"),
        ("Schema.org classes",            "2,987"),
        ("Enrichment cache hit rate",     "~95 %  (in-memory per session)"),
        ("SHACL export time",             "< 100 ms per brick"),
        ("Concurrent sessions tested",    "50+"),
        ("Session memory footprint",      "~10 MB"),
        ("Test suite",                    "60 pytest tests, all passing"),
    ]
    col_l = [Inches(0.5), Inches(7.5)]
    top = Inches(1.9)
    for metric, value in metrics:
        add_textbox(sl, metric, col_l[0], top, Inches(6.8), Inches(0.48),
                    font_size=16, color=C_LIGHT)
        add_textbox(sl, value,  col_l[1], top, Inches(5.5), Inches(0.48),
                    font_size=16, bold=True, color=C_GREEN)
        top += Inches(0.6)


def slide_summary(prs):
    sl = blank_slide(prs)
    fill_bg(sl)
    accent_bar(sl)
    heading(sl, "Key Takeaways", font_size=36)

    items = [
        ("🧱  Bricks",       "Semantic, reusable SHACL components — NodeShape + PropertyShapes as JSON"),
        ("🔍  Enrichment",   "4-layer resolution — generic, not QUDT-only; any ontology works"),
        ("🏗  Multi-tenant", "Web + Qt + API share one BrickCore — no logic duplication"),
        ("🔌  Extensible",   "Add ontologies or TTL rules — zero Python changes required"),
        ("📐  Standards",    "SHACL + DASH output — works with any validator, no vendor lock-in"),
    ]
    top = Inches(1.8)
    for icon_title, desc in items:
        add_textbox(sl, icon_title, Inches(0.5),  top, Inches(3.0), Inches(0.65),
                    font_size=20, bold=True, color=C_ACCENT)
        add_textbox(sl, desc,       Inches(3.7),  top, Inches(9.2), Inches(0.65),
                    font_size=17, color=C_LIGHT)
        top += Inches(0.94)


def slide_qa(prs):
    sl = blank_slide(prs)
    fill_bg(sl)
    accent_bar(sl, C_ORANGE)
    add_textbox(sl, "Questions?",
                Inches(0.5), Inches(1.5), Inches(12.5), Inches(1.4),
                font_size=58, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    qas = [
        ("Why not JSON Schema?",
         "SHACL validates semantics (units, linked data); JSON Schema validates structure. Can convert."),
        ("Why not use SPARQL for ontology queries?",
         "rdflib supports it, but direct traversal is faster, simpler to debug, and sufficient — SPARQL adds value only for complex cross-graph joins."),
        ("Can it work with property graphs (Neo4j)?",
         "SHACL is RDF-native; but SHACL validates before storage — export to JSON/JSON-LD after."),
        ("Why single-file React?",
         "Rapid prototyping, no build step. Production: migrate to Vite + React bundle."),
    ]
    top = Inches(3.3)
    for q, a in qas:
        add_textbox(sl, f"Q: {q}", Inches(0.6), top,              Inches(12.2), Inches(0.42),
                    font_size=15, bold=True, color=C_YELLOW)
        add_textbox(sl, f"A: {a}", Inches(0.6), top + Inches(0.42), Inches(12.2), Inches(0.42),
                    font_size=14, color=C_LIGHT)
        top += Inches(0.98)


# ── Main ──────────────────────────────────────────────────────────────────────

def main():
    prs = new_prs()

    slide_title(prs)           # 1
    slide_agenda(prs)          # 2
    slide_why_shacl(prs)       # 3
    slide_brick_anatomy(prs)   # 4
    slide_architecture(prs)    # 5
    slide_shared_library(prs)  # 6
    slide_enrichment_layer0(prs)    # 7
    slide_enrichment_layers12(prs)  # 8
    slide_enrichment_layers34(prs)  # 9
    slide_widget_rules(prs)    # 10
    slide_shacl_export(prs)    # 11
    slide_dash_mappings(prs)   # 12
    slide_web_frontend(prs)    # 13
    slide_qt_frontend(prs)     # 14
    slide_extension(prs)       # 15
    slide_numbers(prs)         # 16
    slide_summary(prs)         # 17
    slide_qa(prs)              # 18

    out = Path(__file__).parent / "technical_deep_dive_slides.pptx"
    prs.save(str(out))
    print(f"Saved {len(prs.slides)} slides → {out}")


if __name__ == "__main__":
    main()
