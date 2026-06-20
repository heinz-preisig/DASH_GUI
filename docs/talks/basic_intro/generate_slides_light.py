"""
Generate basic_intro PPTX slides for DASH_GUI talk.
Run from repo root:  uv run python docs/talks/basic_intro/generate_slides.py
Output: docs/talks/basic_intro/basic_intro_slides_light.pptx
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

# ── Colour palette (light theme) ─────────────────────────────────────────────
C_BG        = RGBColor(0xFF, 0xFF, 0xFF)   # white background
C_ACCENT    = RGBColor(0x02, 0x6D, 0xB0)   # strong blue
C_WHITE     = RGBColor(0x1A, 0x1A, 0x2E)   # near-black (was white text)
C_LIGHT     = RGBColor(0x33, 0x44, 0x55)   # dark grey body text
C_GREEN     = RGBColor(0x16, 0x7A, 0x40)   # dark green
C_YELLOW    = RGBColor(0x92, 0x5E, 0x00)   # amber (readable on white)
C_RED       = RGBColor(0xBF, 0x1B, 0x1B)   # dark red
C_CODE_BG   = RGBColor(0xF0, 0xF4, 0xF8)   # light grey for code blocks

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs: Presentation):
    blank_layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(blank_layout)


def fill_bg(slide, color: RGBColor = C_BG):
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, text: str, left, top, width, height,
                font_size=24, bold=False, color=C_WHITE,
                align=PP_ALIGN.LEFT, word_wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def add_para(tf, text: str, font_size=20, bold=False,
             color=C_WHITE, align=PP_ALIGN.LEFT, indent_level=0):
    from pptx.util import Pt
    p = tf.add_paragraph()
    p.alignment = align
    p.level = indent_level
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return p


def add_divider(slide, top, color=C_ACCENT, width_fraction=0.85):
    from pptx.util import Pt
    left  = SLIDE_W * (1 - width_fraction) / 2
    width = SLIDE_W * width_fraction
    line  = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.LINE is not directly usable; use connector
        left, top, width, Pt(0)
    )
    line.line.color.rgb = color
    line.line.width     = Pt(1.5)
    return line


def section_pill(slide, text: str, left=None, top=None):
    """Full-width title bar across the top of the slide."""
    bar = slide.shapes.add_shape(1, 0, 0, SLIDE_W, Inches(0.52))
    bar.fill.solid()
    bar.fill.fore_color.rgb = C_ACCENT
    bar.line.fill.background()
    tf = bar.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = f"  {text}"
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)


def add_code_box(slide, code: str, left, top, width, height, font_size=13):
    box = slide.shapes.add_shape(1, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = C_CODE_BG
    box.line.color.rgb = C_ACCENT
    box.line.width = Pt(1)
    tf = box.text_frame
    tf.word_wrap = False
    first = True
    for line in code.split("\n"):
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.color.rgb = RGBColor(0x02, 0x4A, 0x7A)   # dark blue on light grey
        run.font.name = "Courier New"


# ── Individual slide builders ─────────────────────────────────────────────────

def slide_title(prs):
    sl = blank_slide(prs)
    fill_bg(sl)

    # big accent bar on left edge
    bar = sl.shapes.add_shape(1, 0, 0, Inches(0.18), SLIDE_H)
    bar.fill.solid(); bar.fill.fore_color.rgb = C_ACCENT
    bar.line.fill.background()

    add_textbox(sl, "From Bricks to Forms",
                Inches(0.4), Inches(1.6), Inches(12), Inches(1.6),
                font_size=52, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)
    add_textbox(sl, "Building Digital Product Passports with SHACL",
                Inches(0.4), Inches(3.3), Inches(11), Inches(0.8),
                font_size=26, color=C_ACCENT, align=PP_ALIGN.LEFT)
    add_textbox(sl, "A visual approach to structured data",
                Inches(0.4), Inches(4.2), Inches(10), Inches(0.6),
                font_size=20, color=C_LIGHT, align=PP_ALIGN.LEFT)
    add_textbox(sl, "No coding background required",
                Inches(0.4), Inches(5.5), Inches(8), Inches(0.5),
                font_size=16, color=C_YELLOW, align=PP_ALIGN.LEFT)


def slide_problem(prs):
    sl = blank_slide(prs)
    fill_bg(sl)
    section_pill(sl, "The Problem", Inches(0.4), Inches(0.25))
    add_textbox(sl, "Digital Product Passports Are Hard",
                Inches(0.4), Inches(0.75), Inches(12), Inches(0.9),
                font_size=36, bold=True, color=C_WHITE)

    # left column — before
    left_box = sl.shapes.add_shape(1, Inches(0.4), Inches(1.9),
                                   Inches(5.8), Inches(4.4))
    left_box.fill.solid(); left_box.fill.fore_color.rgb = RGBColor(0xFF,0xED,0xED)
    left_box.line.color.rgb = C_RED; left_box.line.width = Pt(1.5)
    tf = left_box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "❌  TODAY"; r.font.size = Pt(18)
    r.font.bold = True; r.font.color.rgb = C_RED
    for line in [
        "• Company A says: 'weight: 5kg'",
        "• Company B says: 'mass: 5000g'",
        "• Company C says: '5'  (just a number!)",
        "",
        "→ Computers can't understand this",
        "→ Manual data entry & errors",
        "→ Compliance headaches",
    ]:
        add_para(tf, line, font_size=17, color=C_LIGHT)

    # right column — after
    right_box = sl.shapes.add_shape(1, Inches(6.8), Inches(1.9),
                                    Inches(5.8), Inches(4.4))
    right_box.fill.solid(); right_box.fill.fore_color.rgb = RGBColor(0xED,0xFF,0xF3)
    right_box.line.color.rgb = C_GREEN; right_box.line.width = Pt(1.5)
    tf2 = right_box.text_frame; tf2.word_wrap = True
    p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = "✓  WITH SHACL"; r2.font.size = Pt(18)
    r2.font.bold = True; r2.font.color.rgb = C_GREEN
    for line in [
        "• Product.weight = 5",
        "• Must be a number",
        "• Must have a unit",
        "• Unit can only be: kg, g, lb, oz",
        "",
        "→ Machines can validate it",
        "→ Forms guide data entry",
        "→ Always consistent",
    ]:
        add_para(tf2, line, font_size=17, color=C_LIGHT)


def slide_vision(prs):
    sl = blank_slide(prs)
    fill_bg(sl)
    section_pill(sl, "The Vision", Inches(0.4), Inches(0.25))
    add_textbox(sl, "What If Data Had Grammar?",
                Inches(0.4), Inches(0.75), Inches(12), Inches(0.9),
                font_size=36, bold=True, color=C_WHITE)

    add_textbox(sl, "Like grammar rules in language, data needs rules too.",
                Inches(0.4), Inches(1.8), Inches(12), Inches(0.6),
                font_size=20, color=C_LIGHT)

    code_bad = '❌  "The product weighs about 5kg usually"   (ambiguous)'
    code_good = (
        "✓   Product.weight  =  5\n"
        "         type  →  number\n"
        "         unit  →  kg  (or g, lb, oz)\n"
        "     required  →  yes\n"
        "      checked  →  automatically"
    )
    add_code_box(sl, code_bad,  Inches(0.7), Inches(2.6), Inches(11.6), Inches(0.75), font_size=16)
    add_code_box(sl, code_good, Inches(0.7), Inches(3.6), Inches(11.6), Inches(2.1),  font_size=17)

    add_textbox(sl, "SHACL is the W3C standard for writing these rules.",
                Inches(0.4), Inches(6.0), Inches(12), Inches(0.55),
                font_size=18, color=C_YELLOW)


def slide_lego(prs):
    sl = blank_slide(prs)
    fill_bg(sl)
    section_pill(sl, "The Brick Idea", Inches(0.4), Inches(0.25))
    add_textbox(sl, "Lego Bricks for Data Schemas",
                Inches(0.4), Inches(0.75), Inches(12), Inches(0.9),
                font_size=36, bold=True, color=C_WHITE)

    add_textbox(sl, "Instead of writing rules from scratch every time → create reusable bricks once, snap them together.",
                Inches(0.4), Inches(1.75), Inches(12.5), Inches(0.7),
                font_size=19, color=C_LIGHT)

    bricks = [
        ("🧱  Mass",        "number + unit\n(kg, g, lb, oz, mg, t)"),
        ("🌡  Temperature", "number + degrees\n(°C, °F, K)"),
        ("⚡  Battery Cap", "number + unit\n(Ah, mAh, Coulombs)"),
        ("🏭  Manufacturer","name + ID + country\n(linked data)"),
    ]
    cols = [Inches(0.4), Inches(3.6), Inches(6.8), Inches(10.0)]
    for i, (title, desc) in enumerate(bricks):
        box = sl.shapes.add_shape(9, cols[i], Inches(2.7), Inches(3.0), Inches(3.4))
        box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0xE8,0xF4,0xFF)
        box.line.color.rgb = C_ACCENT; box.line.width = Pt(1.5)
        tf = box.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = title; r.font.size = Pt(20)
        r.font.bold = True; r.font.color.rgb = C_ACCENT
        add_para(tf, "", font_size=10, color=C_WHITE)
        add_para(tf, desc, font_size=16, color=C_LIGHT, align=PP_ALIGN.CENTER)

    add_textbox(sl, "Define once → reuse across every product schema in your organisation.",
                Inches(0.4), Inches(6.5), Inches(12.5), Inches(0.5),
                font_size=17, color=C_YELLOW)


def slide_brick_editor(prs):
    sl = blank_slide(prs)
    fill_bg(sl)
    section_pill(sl, "Brick Editor", Inches(0.4), Inches(0.25))
    add_textbox(sl, "Creating a Reusable Brick",
                Inches(0.4), Inches(0.75), Inches(12), Inches(0.9),
                font_size=36, bold=True, color=C_WHITE)

    add_textbox(sl, "Point-and-click — no coding needed:",
                Inches(0.4), Inches(1.75), Inches(6), Inches(0.5),
                font_size=19, color=C_LIGHT)

    fields = [
        ("Name",     "BatteryCapacity"),
        ("Path",     "ex:batteryCapacity"),
        ("Datatype", "xsd:decimal"),
        ("Class",    "qudt:ElectricCharge   ← triggers magic ✨"),
    ]
    top = Inches(2.4)
    for label, value in fields:
        row = sl.shapes.add_shape(1, Inches(0.5), top, Inches(6.2), Inches(0.55))
        row.fill.solid(); row.fill.fore_color.rgb = RGBColor(0xEC,0xF4,0xFF)
        row.line.fill.background()
        tf = row.text_frame; tf.word_wrap = False
        p = tf.paragraphs[0]
        r1 = p.add_run(); r1.text = f"{label:<12}"; r1.font.size = Pt(16)
        r1.font.bold = True; r1.font.color.rgb = C_ACCENT; r1.font.name = "Courier New"
        r2 = p.add_run(); r2.text = value; r2.font.size = Pt(16)
        r2.font.color.rgb = C_LIGHT; r2.font.name = "Courier New"
        top += Inches(0.62)

    # right side: result
    add_textbox(sl, "Result →",
                Inches(7.2), Inches(2.1), Inches(2), Inches(0.5),
                font_size=18, bold=True, color=C_YELLOW)
    result_code = (
        "Unit dropdown appears\n"
        "automatically:\n"
        "\n"
        "  ▼  Ampere-hour (Ah)\n"
        "     Milliampere-hour (mAh)\n"
        "     Coulomb (C)\n"
        "     Kiloampere-hour (kAh)"
    )
    add_code_box(sl, result_code, Inches(7.2), Inches(2.7), Inches(5.7), Inches(3.0), font_size=15)

    add_textbox(sl, "The system looks up applicable units from the QUDT ontology — no manual list needed.",
                Inches(0.4), Inches(6.4), Inches(12.5), Inches(0.5),
                font_size=16, color=C_YELLOW)


def slide_schema(prs):
    sl = blank_slide(prs)
    fill_bg(sl)
    section_pill(sl, "Schema Editor", Inches(0.4), Inches(0.25))
    add_textbox(sl, "Assembling Bricks into a Product Schema",
                Inches(0.4), Inches(0.75), Inches(12), Inches(0.9),
                font_size=36, bold=True, color=C_WHITE)

    tree_code = (
        "EV Battery Schema\n"
        "├── Battery Identification\n"
        "│   ├── Model Number      (text)\n"
        "│   ├── Serial Number     (text)\n"
        "│   └── Manufacturer      (linked data)\n"
        "├── Physical Properties\n"
        "│   ├── Mass           ← [Mass brick]        → unit dropdown\n"
        "│   └── Capacity       ← [BatteryCapacity]   → unit dropdown\n"
        "└── Environmental\n"
        "    ├── Carbon Footprint  ← [CO₂ brick]      → unit dropdown\n"
        "    └── Recyclability     ← [Percent brick]  → % field"
    )
    add_code_box(sl, tree_code, Inches(0.4), Inches(1.9), Inches(12.5), Inches(4.5), font_size=15)

    add_textbox(sl, "Visual tree editor — drag bricks in, set connections, done.",
                Inches(0.4), Inches(6.6), Inches(12.5), Inches(0.5),
                font_size=17, color=C_YELLOW)


def slide_magic(prs):
    sl = blank_slide(prs)
    fill_bg(sl)
    section_pill(sl, "One Source → Many Outputs", Inches(0.4), Inches(0.25))
    add_textbox(sl, "The Magic: One Brick, Three Things",
                Inches(0.4), Inches(0.75), Inches(12), Inches(0.9),
                font_size=36, bold=True, color=C_WHITE)

    cols = [
        ("Brick (JSON)", Inches(0.3),
         "{\n  name: 'Mass',\n  datatype: xsd:decimal,\n  class: qudt:Mass,\n  units: [kg,g,lb,...]\n}"),
        ("SHACL Rules", Inches(4.55),
         "sh:property [\n  sh:path ex:mass ;\n  sh:datatype xsd:decimal ;\n  sh:class qudt:Mass ;\n  sh:in (kg g lb oz) ;\n] ."),
        ("User Form", Inches(8.8),
         "  Mass\n  ┌─────────┬──────┐\n  │  5      │ ▼ kg │\n  └─────────┴──────┘\n\n  Change brick once\n  → all forms update"),
    ]
    arrows = [Inches(4.2), Inches(8.45)]
    for ax in arrows:
        add_textbox(sl, "→", ax, Inches(3.8), Inches(0.5), Inches(0.5),
                    font_size=30, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)

    for title, left, code in cols:
        add_textbox(sl, title, left, Inches(1.9), Inches(4.1), Inches(0.5),
                    font_size=17, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
        add_code_box(sl, code, left, Inches(2.5), Inches(4.15), Inches(3.5), font_size=14)


def slide_enrichment(prs):
    sl = blank_slide(prs)
    fill_bg(sl)
    section_pill(sl, "Enrichment Engine", Inches(0.4), Inches(0.25))
    add_textbox(sl, "Smart Widgets from Context",
                Inches(0.4), Inches(0.75), Inches(12), Inches(0.9),
                font_size=36, bold=True, color=C_WHITE)

    steps = [
        ("1", "You type:  qudt:Mass",               C_WHITE),
        ("2", "System queries QUDT ontology",        C_LIGHT),
        ("3", "Finds 6 applicable units",            C_LIGHT),
        ("4", "Creates unit dropdown automatically", C_GREEN),
    ]
    top = Inches(2.0)
    for num, text, col in steps:
        circle = sl.shapes.add_shape(9, Inches(0.5), top, Inches(0.55), Inches(0.55))
        circle.fill.solid(); circle.fill.fore_color.rgb = C_ACCENT
        circle.line.fill.background()
        ctf = circle.text_frame
        cp = ctf.paragraphs[0]; cp.alignment = PP_ALIGN.CENTER
        cr = cp.add_run(); cr.text = num; cr.font.size = Pt(16)
        cr.font.bold = True; cr.font.color.rgb = C_BG
        add_textbox(sl, text, Inches(1.25), top + Inches(0.05),
                    Inches(11), Inches(0.5), font_size=20, color=col)
        if num != "4":
            add_textbox(sl, "↓", Inches(0.62), top + Inches(0.55),
                        Inches(0.4), Inches(0.4), font_size=18,
                        color=C_ACCENT, align=PP_ALIGN.CENTER)
        top += Inches(1.05)

    add_textbox(sl, "Works for any physical quantity: temperature, pressure, energy, voltage …",
                Inches(0.4), Inches(6.4), Inches(12.5), Inches(0.5),
                font_size=17, color=C_YELLOW)


def slide_architecture(prs):
    sl = blank_slide(prs)
    fill_bg(sl)
    section_pill(sl, "Architecture", Inches(0.4), Inches(0.25))
    add_textbox(sl, "How It All Connects",
                Inches(0.4), Inches(0.75), Inches(12), Inches(0.9),
                font_size=36, bold=True, color=C_WHITE)

    layers = [
        ("Shared Ontologies",  "QUDT · Schema.org · FOAF · SKOS  (17 vocabularies)", C_LIGHT,   Inches(1.4)),
        ("Brick Library",      "Reusable building blocks (JSON)",                     C_ACCENT,  Inches(2.35)),
        ("Schema Composer",    "Assemble bricks into product blueprints",             C_ACCENT,  Inches(3.3)),
        ("SHACL Exporter",     "Generates validation rules (Turtle)",                 C_LIGHT,   Inches(4.25)),
        ("Form Generator",     "HTML interface the data-entry clerk sees",            C_GREEN,   Inches(5.2)),
        ("Validated Data",     "Ready for Digital Product Passports",                 C_YELLOW,  Inches(6.15)),
    ]
    for label, desc, col, top in layers:
        box = sl.shapes.add_shape(1, Inches(1.5), top, Inches(10), Inches(0.72))
        box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0xF5,0xF8,0xFF)
        box.line.color.rgb = col; box.line.width = Pt(1.5)
        tf = box.text_frame; tf.word_wrap = False
        p = tf.paragraphs[0]
        r1 = p.add_run(); r1.text = f"{label:<22}"; r1.font.size = Pt(17)
        r1.font.bold = True; r1.font.color.rgb = col; r1.font.name = "Courier New"
        r2 = p.add_run(); r2.text = desc; r2.font.size = Pt(15)
        r2.font.color.rgb = C_LIGHT; r2.font.name = "Courier New"


def slide_impact(prs):
    sl = blank_slide(prs)
    fill_bg(sl)
    section_pill(sl, "Real-World Impact", Inches(0.4), Inches(0.25))
    add_textbox(sl, "Why This Matters",
                Inches(0.4), Inches(0.75), Inches(12), Inches(0.9),
                font_size=36, bold=True, color=C_WHITE)

    rows = [
        ("Schema definition time", "3 months",           "1 week",       ),
        ("Unit configuration",     "manual lookup",      "automatic",    ),
        ("Validation errors",      "50+ (found late)",   "2 (caught early)"),
        ("Consistency",            "varies per team",    "standardised", ),
        ("Non-technical authors",  "impossible",         "point-and-click"),
    ]
    headers = ("", "Before", "After")
    col_l = [Inches(0.4), Inches(5.2), Inches(9.3)]
    top = Inches(2.0)
    # header row
    for i, h in enumerate(headers):
        add_textbox(sl, h, col_l[i], top, Inches(4.5), Inches(0.45),
                    font_size=17, bold=True,
                    color=C_ACCENT if i > 0 else C_LIGHT)
    top += Inches(0.5)
    for metric, before, after in rows:
        add_textbox(sl, metric, col_l[0], top, Inches(4.6), Inches(0.45),
                    font_size=16, color=C_LIGHT)
        add_textbox(sl, before,  col_l[1], top, Inches(3.8), Inches(0.45),
                    font_size=16, color=C_RED)
        add_textbox(sl, after,   col_l[2], top, Inches(3.8), Inches(0.45),
                    font_size=16, color=C_GREEN)
        top += Inches(0.72)


def slide_summary(prs):
    sl = blank_slide(prs)
    fill_bg(sl)

    bar = sl.shapes.add_shape(1, 0, 0, Inches(0.18), SLIDE_H)
    bar.fill.solid(); bar.fill.fore_color.rgb = C_ACCENT
    bar.line.fill.background()

    add_textbox(sl, "Key Takeaways",
                Inches(0.4), Inches(0.6), Inches(12), Inches(0.9),
                font_size=36, bold=True, color=C_WHITE)

    items = [
        ("🧱  Bricks",   "Reusable, semantic data components — define once, use everywhere"),
        ("🔗  SHACL",    "Standardised validation rules — machines can check your data"),
        ("✨  Enrichment","Ontology-driven smart widgets — unit dropdowns appear automatically"),
        ("🎨  Forms",    "Automatic user interfaces — no coding needed for data entry"),
        ("♻️  Reuse",    "Change a brick once → every form using it updates"),
    ]
    top = Inches(1.8)
    for icon_title, desc in items:
        add_textbox(sl, icon_title, Inches(0.5), top, Inches(3.2), Inches(0.65),
                    font_size=20, bold=True, color=C_ACCENT)
        add_textbox(sl, desc, Inches(3.8), top, Inches(9.0), Inches(0.65),
                    font_size=18, color=C_LIGHT)
        top += Inches(0.92)


def slide_getting_started(prs):
    sl = blank_slide(prs)
    fill_bg(sl)
    section_pill(sl, "Getting Started")
    add_textbox(sl, "One Download, One Command",
                Inches(0.4), Inches(0.65), Inches(12.5), Inches(0.9),
                font_size=36, bold=True, color=C_WHITE)

    steps = [
        ("1", "Download",  "Clone or download the repository from GitHub"),
        ("2", "Run",       "Execute the included start script:  dev-start-both-docker.sh"),
        ("3", "Open",      "Browser opens automatically — no installation needed"),
        ("4", "Your data", "Library folder stays on your machine — nothing is sent anywhere"),
    ]
    top = Inches(1.85)
    for num, label, desc in steps:
        circle = sl.shapes.add_shape(9, Inches(0.4), top, Inches(0.55), Inches(0.55))
        circle.fill.solid(); circle.fill.fore_color.rgb = C_ACCENT
        circle.line.fill.background()
        ctf = circle.text_frame
        cp = ctf.paragraphs[0]; cp.alignment = PP_ALIGN.CENTER
        cr = cp.add_run(); cr.text = num; cr.font.size = Pt(16)
        cr.font.bold = True; cr.font.color.rgb = RGBColor(0xFF,0xFF,0xFF)
        add_textbox(sl, label, Inches(1.15), top + Inches(0.04), Inches(1.8), Inches(0.48),
                    font_size=18, bold=True, color=C_ACCENT)
        add_textbox(sl, desc,  Inches(3.1),  top + Inches(0.07), Inches(9.8), Inches(0.44),
                    font_size=17, color=C_LIGHT)
        top += Inches(1.05)

    add_textbox(sl, "The script handles everything — Docker is just the packaging, you never need to know it's there.",
                Inches(0.4), Inches(6.4), Inches(12.5), Inches(0.55),
                font_size=15, color=C_YELLOW)


def slide_qa(prs):
    sl = blank_slide(prs)
    fill_bg(sl)

    bar = sl.shapes.add_shape(1, 0, 0, Inches(0.18), SLIDE_H)
    bar.fill.solid(); bar.fill.fore_color.rgb = C_YELLOW
    bar.line.fill.background()

    add_textbox(sl, "Questions?",
                Inches(0.4), Inches(1.8), Inches(12), Inches(1.4),
                font_size=60, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    qas = [
        ("Can non-programmers use this?",
         "Yes — brick_app and schema_app are point-and-click web tools."),
        ("What about our own vocabulary?",
         "Load any RDF/OWL ontology — enrichment works with it automatically."),
        ("Does it connect to our database?",
         "Exported SHACL validates data; can export to JSON/XML/RDF for any system."),
    ]
    top = Inches(3.5)
    for q, a in qas:
        add_textbox(sl, f"Q: {q}", Inches(0.6), top, Inches(12.2), Inches(0.42),
                    font_size=16, bold=True, color=C_YELLOW)
        add_textbox(sl, f"A: {a}", Inches(0.6), top + Inches(0.42), Inches(12.2), Inches(0.42),
                    font_size=15, color=C_LIGHT)
        top += Inches(1.05)


# ── Main ──────────────────────────────────────────────────────────────────────

def main():
    prs = new_prs()

    slide_title(prs)
    slide_problem(prs)
    slide_vision(prs)
    slide_lego(prs)
    slide_brick_editor(prs)
    slide_schema(prs)
    slide_magic(prs)
    slide_enrichment(prs)
    slide_architecture(prs)
    slide_impact(prs)
    slide_summary(prs)
    slide_getting_started(prs)
    slide_qa(prs)

    out = Path(__file__).parent / "basic_intro_slides_light.pptx"
    prs.save(str(out))
    print(f"Saved {len(prs.slides)} slides → {out}")


if __name__ == "__main__":
    main()
